#!/usr/bin/env python3
"""
make_deck.py — builds a branded PowerPoint (and, with LibreOffice, a PDF) from data.json.

Weekly data ONLY (live web sections are intentionally excluded from the deck).

Usage:
    python make_deck.py data.json            # -> DCM_Weekly_<week>.pptx
    python make_deck.py data.json --pdf      # also writes .pdf via LibreOffice
"""
import sys, json, os, subprocess, datetime
# ensure imports (commentary.py) resolve from this script's own folder
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---- brand palette ----
NAVY  = RGBColor(0x1E,0x27,0x61)
DEEP  = RGBColor(0x12,0x16,0x3A)
GOLD  = RGBColor(0xC9,0xA2,0x27)
ICE   = RGBColor(0xCA,0xDC,0xFC)
ICEBG = RGBColor(0xEE,0xF3,0xFC)
INK   = RGBColor(0x23,0x27,0x3B)
MUTED = RGBColor(0x6B,0x72,0x80)
WHITE = RGBColor(0xFF,0xFF,0xFF)
GOOD  = RGBColor(0x1E,0x8E,0x5A)
BAD   = RGBColor(0xC0,0x39,0x2B)
LTGREY= RGBColor(0xB9,0xC3,0xE4)
HEAD_FONT="Cambria"      # safe serif
BODY_FONT="Calibri"      # safe sans

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FOOT=""       # set per-build to "Bangladesh Macro Weekly  -  week ending <date>"
PAGE=[1]      # running page counter (mutable)

def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb=color; shape.line.fill.background()

def box(slide,x,y,w,h,fill=None,line=None,shadow=False,radius=False):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                               Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is not None: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    else: shp.fill.background()
    if line is not None: shp.line.color.rgb=line; shp.line.width=Pt(0.75)
    else: shp.line.fill.background()
    shp.shadow.inherit=False
    if shadow:
        el=shp._element.spPr
    return shp

def text(slide,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,
         font=BODY_FONT,wrap=True,space=None):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    tf.vertical_anchor=anchor
    for m in (tf.margin_left,tf.margin_right): pass
    tf.margin_left=Pt(2);tf.margin_right=Pt(2);tf.margin_top=Pt(1);tf.margin_bottom=Pt(1)
    if isinstance(runs,str): runs=[(runs,{})]
    first=True
    for line in runs if isinstance(runs[0],list) else [runs]:
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        first=False; p.alignment=align
        if space: p.space_after=Pt(space)
        for t,st in (line if isinstance(line,list) else [line]):
            r=p.add_run(); r.text=t
            r.font.name=st.get("font",font); r.font.size=Pt(st.get("sz",14))
            r.font.bold=st.get("b",False); r.font.italic=st.get("i",False)
            r.font.color.rgb=st.get("c",INK)
    return tb

def R(t,**st): return (t,st)

# ---------------- deck ----------------
def build(data, out_pptx):
    prs=Presentation(); prs.slide_width=EMU_W; prs.slide_height=EMU_H
    blank=prs.slide_layouts[6]
    meta=data["meta"]; kpis={k["ind"]:k for k in data["kpis"]}; S=data.get("series",{})
    # resolved commentary (manual stacked with auto) — shared with website
    CR = data.get("commentaryResolved")
    if not CR:
        from commentary import resolve_section
        CR = {sec: resolve_section(sec, (data.get("commentary") or {}).get(sec,""), kpis)
              for sec in ["fx","global","rates","ycurve","inflation","external","reserves","commodity"]}
    def takeaways(sec):
        return [(l["text"], bool(l["manual"])) for l in (CR.get(sec) or [])]
    global FOOT
    FOOT=f"Bangladesh Macro Weekly  -  week ending {meta.get('weekEnding','')}"
    PAGE[0]=1
    def g(ind,f="cur"): 
        k=kpis.get(ind,{}); return k.get(f)
    def fmtn(v,d=2): 
        return "—" if v is None else f"{v:,.{d}f}"

    # ---------- SLIDE 1 : title ----------
    s=prs.slides.add_slide(blank)
    box(s,0,0,13.333,7.5,fill=DEEP)
    box(s,0,0,13.333,0.12,fill=GOLD)
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.9),Inches(1.7),Inches(0.22),Inches(0.22)); solid(dot,GOLD)
    text(s,1.2,1.55,9.0,0.5,[[R("BANGLADESH  MACRO",font=BODY_FONT,sz=15,b=True,c=WHITE)]])
    text(s,0.85,2.2,11.8,1.6,[[R(meta.get("title","DCM Weekly : দেশী Macro"),font=HEAD_FONT,sz=46,b=True,c=WHITE)]])
    text(s,0.9,3.65,11.4,0.5,[[R(meta.get("subtitle",""),font=HEAD_FONT,sz=20,i=True,c=ICE)]])
    # headline box (auto or manual intro paragraph)
    box(s,0.9,4.5,11.5,1.6,fill=RGBColor(0x1A,0x20,0x52),radius=True)
    text(s,1.2,4.66,10.9,1.3,[[R(meta.get("headline",""),sz=14,c=ICE)]])
    text(s,0.9,6.45,7.0,0.4,[[R("Week ending  ",sz=13,c=ICE),R(meta.get("weekEnding","—"),sz=13,b=True,c=WHITE)]])
    text(s,0.9,6.85,11.0,0.4,[[R(meta.get("forAudience",""),sz=11,c=MUTED)]])
    text(s,9.5,6.45,2.9,0.4,[[R("Prepared by",sz=11,c=MUTED)]],align=PP_ALIGN.RIGHT)
    text(s,9.5,6.72,2.9,0.4,[[R(meta.get("preparedBy",""),sz=11,b=True,c=ICE)]],align=PP_ALIGN.RIGHT)

    # ---------- SLIDE 2 : executive dashboard (stat grid) ----------
    s=prs.slides.add_slide(blank); base(s,"Executive Dashboard","AT A GLANCE",FOOT,2)
    text(s,0.5,1.55,11.8,0.34,[[R("Week-on-week moves. Green = favourable, red = adverse, grey = market/neutral.",sz=11,i=True,c=MUTED)]])
    dash=[k for k in data["kpis"] if k.get("dash")][:16]
    cols=4; cw=2.9; ch=1.11; gx=0.25; gy=0.19; x0=0.5; y0=1.95
    for i,k in enumerate(dash):
        r,c=divmod(i,cols); x=x0+c*(cw+gx); y=y0+r*(ch+gy)
        card=box(s,x,y,cw,ch,fill=ICEBG,line=RGBColor(0xDD,0xE2,0xEC),radius=True,shadow=True)
        # status dot
        cur,prev=k.get("cur"),k.get("prev"); good=k.get("good")
        if cur is not None and prev is not None and abs(cur-prev)>1e-9 and good is not None:
            dotc=GOOD if (cur-prev>0)==(good==1) else BAD
        elif good is None: dotc=GOLD
        else: dotc=NAVY
        dd=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.18),Inches(y+0.2),Inches(0.1),Inches(0.1)); solid(dd,dotc)
        unit=k.get("unit") or ""
        val=fmtn(k.get("cur"),k.get("dec",2))+("%" if unit=="%" else (" "+unit if unit else ""))
        text(s,x+0.34,y+0.12,cw-0.45,0.3,[[R(k["ind"],sz=11,b=True,c=MUTED)]])
        text(s,x+0.18,y+0.4,cw-0.3,0.42,[[R(val,font=HEAD_FONT,sz=22,b=True,c=NAVY)]])
        if cur is not None and prev is not None:
            d=cur-prev
            col=MUTED if (good is None or abs(d)<1e-9) else (GOOD if (d>0)==(good==1) else BAD)
            ar="▲" if d>1e-9 else ("▼" if d<-1e-9 else "-")
            sfx="%" if k.get("chg")=="pct" else (" pp" if unit=="%" else "")
            shown=(d/abs(prev)*100) if k.get("chg")=="pct" and prev else d
            text(s,x+0.18,y+0.79,cw-0.3,0.3,[[R(f"{ar} {abs(shown):,.2f}{sfx}  WoW",sz=10,b=True,c=col)]])
    PAGE[0]=3

    # ---------- content slides with native charts (reference layout) ----------
    lab=[d[2:7] for d in S.get("dates",[])]

    chart_slide(prs,blank,"FX & Exchange Rates","CURRENCY",
        "Exchange rate: interbank (WAR) vs BC selling",
        lab,"line",[("Interbank (WAR)",S.get("Interbank")),("BC Selling",S.get("BC_Selling"))],
        takeaways("fx"),kpis,[("Interbank (WAR)","Interbank (WAR)"),("BC Selling (Prime)","BC Selling (Prime)")])

    global_slide(prs,blank,kpis,takeaways("global"))

    chart_slide(prs,blank,"Domestic Rates","POLICY & MONEY MARKETS",
        "Call money vs policy repo rate (%)",
        lab,"line",[("Call money",S.get("CallMoney")),("Policy repo",S.get("Policy"))],
        takeaways("rates"),kpis,[("Policy (Repo) Rate","Policy (Repo) Rate"),("Bank Rate","Bank Rate"),("Call Money","Call Money")])

    yc_slide(prs,blank,kpis,takeaways("ycurve"))

    chart_slide(prs,blank,"Inflation","PRICES",
        "Inflation: point-to-point vs 12-month average (%)",
        lab,"line",[("Point-to-point",S.get("P2P_Infl")),("12-month avg",S.get("Avg12M_Infl"))],
        takeaways("inflation"),kpis,[("P2P Inflation","P2P Inflation"),("Inflation (12M avg)","Inflation (12M avg)")])

    chart_slide(prs,blank,"Trade & Remittances","INFLOWS & TRADE",
        "Trade & remittances, monthly ($bn)",
        lab,"line",[("Export",S.get("Export")),("Import",S.get("Import")),("Remittance",S.get("Remittance"))],
        takeaways("external"),kpis,[("Import (monthly)","Import (monthly)"),("Export (monthly)","Export (monthly)"),("Remittance (monthly)","Remittance (monthly)")])

    chart_slide(prs,blank,"Foreign Exchange Reserves","EXTERNAL BUFFER",
        "Foreign exchange reserves ($bn)",
        lab,"line",[("Gross",S.get("Reserve_Gross")),("BPM6",S.get("Reserve_BPM6"))],
        takeaways("reserves"),kpis,[("Reserve Gross","Reserve Gross"),("Reserve BPM6","Reserve BPM6")])

    commodity_slide(prs,blank,lab,S,kpis,takeaways("commodity"))

    # ---------- special issue ----------
    sp=data.get("special",{})
    if sp.get("title"):
        special_slide(prs,blank,sp)
    # ---------- data calendar ----------
    cal=data.get("calendar",[])
    if cal:
        s=prs.slides.add_slide(blank); base(s,"Data Calendar & Watch","LOOKING AHEAD",FOOT,PAGE[0]); PAGE[0]+=1
        y=1.8
        for c in cal:
            box(s,0.5,y,12.33,0.62,fill=ICEBG if (cal.index(c)%2==0) else WHITE,
                line=RGBColor(0xDD,0xE2,0xEC),radius=True,shadow=True)
            text(s,0.75,y+0.13,2.6,0.4,[[R(str(c.get("date","")),sz=13,b=True,c=GOLD)]])
            text(s,3.4,y+0.13,9.2,0.4,[[R(str(c.get("event","")),sz=13,c=INK)]])
            y+=0.72

    # ---------- closing / sources ----------
    s=prs.slides.add_slide(blank); box(s,0,0,13.333,7.5,fill=DEEP); box(s,0,0,13.333,0.12,fill=GOLD)
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.9),Inches(1.32),Inches(0.16),Inches(0.16)); solid(dot,GOLD)
    text(s,1.14,1.2,9.0,0.34,[[R("METHODOLOGY & SOURCES",sz=12,b=True,c=WHITE)]])
    text(s,0.9,1.75,11.5,0.9,[[R("Sources & Disclaimer",font=HEAD_FONT,sz=30,b=True,c=WHITE)]])
    text(s,0.9,2.9,11.5,1.4,[
        [R("Primary sources: Bangladesh Bank (FX, reserves, policy & market rates, remittances); "
          "BBS (CPI); EPB (trade); Dhaka Stock Exchange. Global: CME Group (Term SOFR), EMMI (EURIBOR), "
          "LBMA/ICE (gold, Brent). Figures are point-in-time and subject to revision by the issuing source.",sz=13,c=ICE)]])
    box(s,0.9,4.4,11.5,1.6,fill=RGBColor(0x1A,0x20,0x52),radius=True)
    text(s,1.2,4.56,10.9,1.3,[[R("Prepared for internal senior-management information only. Not investment advice, "
        "an offer, or a solicitation. Weekly figures are compiled from public releases and may be provisional. "
        "Verify against primary sources before acting. Distribution is restricted and confidential.",sz=12,c=RGBColor(0x9A,0xA3,0xC8))]])
    text(s,0.9,6.7,11.5,0.4,[[R(FOOT,sz=9,c=MUTED)]])

    prs.save(out_pptx)
    return out_pptx

def base(slide,title,eyebrow,footer=None,page=None):
    box(slide,0,0,13.333,7.5,fill=WHITE)
    dot=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.5),Inches(0.62),Inches(0.16),Inches(0.16))
    solid(dot,GOLD)
    text(slide,0.74,0.50,9.0,0.34,[[R(eyebrow,sz=12,b=True,c=GOLD)]])
    text(slide,0.48,0.84,12.2,0.7,[[R(title,font=HEAD_FONT,sz=30,b=True,c=NAVY)]])
    if footer:
        text(slide,0.5,7.06,6.0,0.3,[[R(footer,sz=9,c=MUTED)]])
        text(slide,11.0,7.06,1.83,0.3,[[R("Internal Usage Only",sz=9,c=MUTED)]],align=PP_ALIGN.RIGHT)
        if page is not None:
            text(slide,6.4,7.06,0.5,0.3,[[R(str(page),sz=9,c=MUTED)]],align=PP_ALIGN.CENTER)

def _chart_frame(slide,title):
    """White rounded card holding a chart, with an italic-ish bold title inside."""
    box(slide,0.5,1.75,7.1,4.0,fill=WHITE,line=RGBColor(0xDD,0xE2,0xEC),radius=True,shadow=True)
    text(slide,0.72,1.86,6.7,0.3,[[R(title,sz=12,b=True,c=NAVY)]])

def add_native_chart(slide,x,y,w,h,cats,series,kind="line"):
    cd=CategoryChartData(); cd.categories=cats
    for name,vals in series:
        cd.add_series(name,[ (None if v is None else v) for v in (vals or []) ])
    ctype=XL_CHART_TYPE.COLUMN_CLUSTERED if kind=="bar" else XL_CHART_TYPE.LINE_MARKERS
    gf=slide.shapes.add_chart(ctype,Inches(x),Inches(y),Inches(w),Inches(h),cd)
    ch=gf.chart
    ch.has_title=False
    ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False
    ch.legend.font.size=Pt(11); ch.legend.font.name=BODY_FONT
    palette=[NAVY,GOLD,GOOD]
    if kind=="bar":
        for i,plot in enumerate(ch.plots):
            plot.gap_width=80
        for i,se in enumerate(ch.series):
            se.format.fill.solid(); se.format.fill.fore_color.rgb=palette[i%len(palette)]
    else:
        for i,se in enumerate(ch.series):
            col=palette[i%len(palette)]
            se.format.line.color.rgb=col; se.format.line.width=Pt(2.25); se.smooth=False
            try:
                m=se.marker; m.style=8
                m.format.fill.solid(); m.format.fill.fore_color.rgb=col; m.format.line.color.rgb=col
            except Exception: pass
    ca=ch.category_axis; va=ch.value_axis
    for ax in (ca,va):
        ax.tick_labels.font.size=Pt(10); ax.tick_labels.font.name=BODY_FONT
        ax.tick_labels.font.color.rgb=MUTED
    va.has_major_gridlines=True; ca.has_major_gridlines=False
    return gf

def stat_callouts(s,pairs,kpis,x0=7.95,y=1.85):
    """Up to 2 per row: label + big value (reference style, no card)."""
    positions=[(x0,y),(10.30,y),(x0,y+0.95),(10.30,y+0.95)]
    for i,(ind,label) in enumerate(pairs):
        if i>=len(positions): break
        px,py=positions[i]; k=kpis.get(ind,{})
        unit=k.get("unit") or ""
        val=("—" if k.get("cur") is None else f"{k.get('cur'):,.{k.get('dec',2)}f}")
        val += ("%" if unit=="%" else (" "+unit if unit else ""))
        text(s,px,py,2.25,0.28,[[R(label,sz=11,b=True,c=MUTED)]])
        text(s,px,py+0.24,2.25,0.4,[[R(val,font=HEAD_FONT,sz=21,b=True,c=NAVY)]])
    # returns y just below the callouts
    rows=(len(pairs)+1)//2
    return y+0.95*rows+0.05

def takeaway_block(s,lines,y,x=7.95,w=4.9):
    text(s,x,y,w,0.3,[[R("KEY TAKEAWAYS",sz=12,b=True,c=GOLD)]])
    yy=y+0.34
    for txt,is_manual in lines:
        mark="◆" if is_manual else "▪"
        col=NAVY if is_manual else INK
        # marker
        text(s,x,yy,0.3,0.4,[[R(mark,sz=11,b=is_manual,c=GOLD)]])
        # wrapped text
        tb=text(s,x+0.3,yy,w-0.3,0.8,[[R(txt,sz=12.5,b=is_manual,c=col)]])
        # estimate height by char count
        lines_est=max(1,int(len(txt)/58)+1)
        yy+=0.30*lines_est+0.06
    return yy

def chart_slide(prs,blank,title,eyebrow,chart_title,cats,kind,series,takeaways,kpis,callouts,foot=None):
    foot=foot or FOOT
    s=prs.slides.add_slide(blank); base(s,title,eyebrow,foot,PAGE[0]); PAGE[0]+=1
    _chart_frame(s,chart_title)
    add_native_chart(s,0.62,2.25,6.9,3.35,cats,series,kind)
    y=stat_callouts(s,callouts,kpis)
    takeaway_block(s,takeaways,y+0.15)

def global_slide(prs,blank,kpis,takeaways):
    s=prs.slides.add_slide(blank); base(s,"Global Reference Rates","GLOBAL FUNDING",FOOT,PAGE[0]); PAGE[0]+=1
    _chart_frame(s,"Global reference rates by tenor (%)")
    g=lambda n:kpis.get(n,{}).get("cur")
    add_native_chart(s,0.62,2.30,6.9,3.30,["3M","6M","12M"],
        [("SOFR",[g("SOFR 3M"),g("SOFR 6M"),g("SOFR 12M")]),
         ("EURIBOR",[g("EURIBOR 3M"),g("EURIBOR 6M"),g("EURIBOR 12M")])],"bar")
    takeaway_block(s,takeaways,1.85)

def yc_slide(prs,blank,kpis,takeaways):
    s=prs.slides.add_slide(blank); base(s,"Yield Curve","GOVERNMENT SECURITIES",FOOT,PAGE[0]); PAGE[0]+=1
    _chart_frame(s,"Government securities yield curve (%)")
    tenors=[("T-Bill 91D","91D"),("T-Bill 182D","182D"),("T-Bill 364D","364D"),
            ("T-Bond 2Y","2Y"),("T-Bond 5Y","5Y"),("T-Bond 10Y","10Y"),("T-Bond 20Y","20Y")]
    cats=[l for _,l in tenors]
    cur=[kpis.get(t,{}).get("cur") for t,_ in tenors]
    prev=[kpis.get(t,{}).get("prev") for t,_ in tenors]
    # synthetic 'year earlier' offset like reference (prev-ish shape)
    add_native_chart(s,0.62,2.25,6.9,3.35,cats,[("This week",cur),("Last week",prev)],"line")
    # yields grid 2-col
    text(s,7.95,1.85,4.9,0.3,[[R("YIELDS (%)",sz=12,b=True,c=GOLD)]])
    grid=[("91D","T-Bill 91D"),("182D","T-Bill 182D"),("364D","T-Bill 364D"),
          ("2Y","T-Bond 2Y"),("5Y","T-Bond 5Y"),("10Y","T-Bond 10Y"),("20Y","T-Bond 20Y")]
    y0=2.19
    for i,(lbl,ind) in enumerate(grid):
        col=i%2; row=i//2
        px=7.95 if col==0 else 10.40
        py=y0+row*0.42
        v=kpis.get(ind,{}).get("cur")
        text(s,px,py,1.0,0.3,[[R(lbl,sz=12,b=True,c=NAVY)]])
        text(s,px+0.9,py,1.2,0.3,[[R("—" if v is None else f"{v:.2f}",sz=12,b=True,c=INK)]])
    takeaway_block(s,takeaways,3.97)

def commodity_slide(prs,blank,lab,S,kpis,takeaways):
    s=prs.slides.add_slide(blank); base(s,"Commodities","GLOBAL MARKETS",FOOT,PAGE[0]); PAGE[0]+=1
    box(s,0.5,1.75,7.1,4.0,fill=WHITE,line=RGBColor(0xDD,0xE2,0xEC),radius=True,shadow=True)
    text(s,0.72,1.86,3.2,0.3,[[R("Gold ($/oz)",sz=12,b=True,c=NAVY)]])
    text(s,4.0,1.86,3.2,0.3,[[R("Brent crude ($/bbl)",sz=12,b=True,c=NAVY)]])
    add_native_chart(s,0.60,2.30,3.30,3.25,lab,[("Gold",S.get("Gold"))],"line")
    add_native_chart(s,3.95,2.30,3.40,3.25,lab,[("Brent",S.get("Oil"))],"line")
    stat_callouts(s,[("Gold","Gold"),("Oil (Brent)","Oil (Brent)")],kpis)
    takeaway_block(s,takeaways,2.93)

def special_slide(prs,blank,sp):
    s=prs.slides.add_slide(blank); box(s,0,0,13.333,7.5,fill=WHITE)
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.5),Inches(0.62),Inches(0.16),Inches(0.16)); solid(dot,GOLD)
    tag=" · ".join([x for x in ["SPECIAL ISSUE",sp.get("tag"),sp.get("date")] if x])
    text(s,0.74,0.50,9.0,0.34,[[R(tag.upper(),sz=12,b=True,c=GOLD)]])
    text(s,0.48,0.84,12.2,0.8,[[R(sp["title"],font=HEAD_FONT,sz=28,b=True,c=NAVY)]])
    # navy briefing card with gold left edge
    card=box(s,0.5,1.95,12.33,4.6,fill=DEEP,radius=True)
    edge=box(s,0.5,1.95,0.12,4.6,fill=GOLD)
    text(s,0.9,2.2,11.0,0.3,[[R("ANALYST BRIEFING",sz=11,b=True,c=GOLD)]])
    y=2.75
    for p in sp.get("points",[]):
        text(s,0.9,y,0.3,0.4,[[R("◆",sz=12,c=GOLD)]])
        text(s,1.25,y,11.2,0.8,[[R(p,sz=14,c=ICE)]])
        y+=0.30*(int(len(p)/95)+1)+0.28
    text(s,0.5,7.06,6.0,0.3,[[R(FOOT,sz=9,c=MUTED)]])
    text(s,11.0,7.06,1.83,0.3,[[R("Internal Usage Only",sz=9,c=MUTED)]],align=PP_ALIGN.RIGHT)

def _find_soffice():
    """Locate a LibreOffice/soffice executable across platforms."""
    import shutil
    # 1) on PATH (Linux/macOS, or Windows if user added it)
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return [p]
    # 2) common Windows install locations
    for p in (r"C:\Program Files\LibreOffice\program\soffice.exe",
              r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"):
        if os.path.exists(p):
            return [p]
    # 3) common macOS location
    mac = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.exists(mac):
        return [mac]
    # 4) the skill helper (this container only)
    helper = "/mnt/skills/public/pptx/scripts/office/soffice.py"
    if os.path.exists(helper):
        return [sys.executable, helper]
    return None

def to_pdf(pptx_path):
    exe = _find_soffice()
    if not exe:
        print("  (PDF step skipped: LibreOffice not found. Install it from libreoffice.org,")
        print("   or open the .pptx and 'Save As PDF'. The PowerPoint was still created.)")
        return None
    outdir = os.path.dirname(os.path.abspath(pptx_path)) or "."
    cmd = exe + ["--headless", "--convert-to", "pdf", "--outdir", outdir, pptx_path]
    try:
        subprocess.run(cmd, check=True, timeout=240)
    except Exception as e:
        print("  (PDF step failed:", e, "- the .pptx is fine; export PDF manually if needed.)")
        return None
    return os.path.splitext(pptx_path)[0] + ".pdf"

def main():
    args=[a for a in sys.argv[1:] if not a.startswith("--")]
    src=args[0] if args else "data.json"
    src=os.path.abspath(src)
    data=json.load(open(src,encoding="utf-8"))
    wk=(data["meta"].get("weekEnding","week") or "week").replace(",","").replace(" ","_")
    outdir=os.path.dirname(src) or "."
    out=os.path.join(outdir, f"DCM_Weekly_{wk}.pptx")
    build(data,out); print("Wrote",out)
    if "--pdf" in sys.argv:
        pdf=to_pdf(out); print("Wrote",pdf)

if __name__=="__main__":
    main()
